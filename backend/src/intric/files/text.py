from enum import Enum
from pathlib import Path

import magic
import pptx
from docx2python import docx2python
from pypdf import PdfReader


class MimeTypesBase(str, Enum):
    @classmethod
    def has_value(cls, value) -> bool:
        base_value = value.split(";")[0].strip()
        return any(base_value == item.value for item in cls)

    @classmethod
    def values(cls):
        return [item.value for item in cls]


class TextMimeTypes(MimeTypesBase):
    MD = "text/markdown"
    TXT = "text/plain"
    PDF = "application/pdf"
    DOCX = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    TEXT_CSV = "text/csv"
    APP_CSV = "application/csv"
    PPTX = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


class TextSanitizer:
    @staticmethod
    def sanitize(text: str) -> str:
        text = text.replace("\x00", "")
        return text


class TextExtractor:
    @staticmethod
    def extract_from_plain_text(filepath: Path) -> str:
        return filepath.read_text("utf-8")

    @staticmethod
    def extract_from_pdf(filepath: Path) -> str:
        reader = PdfReader(filepath)
        extracted_text = " ".join([page.extract_text() for page in reader.pages])
        sanitized_text = TextSanitizer.sanitize(extracted_text)
        return sanitized_text

    @staticmethod
    def extract_from_docx(filepath: Path) -> str:
        with docx2python(filepath) as docx_content:
            return docx_content.text

    @staticmethod
    def extract_from_pptx(filepath: Path) -> str:
        # Extract text from pptx using python-pptx
        extracted_text = ""
        presentation = pptx.Presentation(filepath)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            extracted_text += run.text + " "
                    extracted_text += "\n"

        return extracted_text

    def extract(self, filepath: Path, mimetype: str | None = None) -> str:
        mimetype = mimetype or magic.from_file(filepath, mime=True)

        match mimetype:
            case (
                TextMimeTypes.TXT
                | TextMimeTypes.MD
                | TextMimeTypes.TEXT_CSV
                | TextMimeTypes.APP_CSV
            ):
                extracted_text = self.extract_from_plain_text(filepath)
            case TextMimeTypes.PDF:
                extracted_text = self.extract_from_pdf(filepath)
            case TextMimeTypes.DOCX:
                extracted_text = self.extract_from_docx(filepath)
            case TextMimeTypes.PPTX:
                extracted_text = self.extract_from_pptx(filepath)
            case _:
                # Fallback to plain text
                extracted_text = self.extract_from_plain_text(filepath)

        return extracted_text.strip()
